# Supported file types:
# - Text: .txt, .csv, .json
# - PDFs: .pdf
# - Word: .docx (modern), .doc (97-2003)
# - PowerPoint: .pptx, .ppt
# - Excel: .xlsx, .xls
# - Images: .jpg, .jpeg, .heic, .png, .tiff, tif, .bmp, .gif (OCR)
# - Emails: .msg (Outlook), .eml
# - Archives: .zip, .7z, .rar (with unpacking)

import os
import csv
import zipfile
import shutil
import tempfile
import spacy
from presidio_analyzer import AnalyzerEngine
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import olefile
import win32com.client
import extract_msg
import email
from email import policy
from email.parser import BytesParser
import pandas as pd
from pillow_heif import register_heif_opener

# Initialize the Analyzer Engine
analyzer = AnalyzerEngine()

# Folder to scan
folder_path = r"N:\Sample Test PII"
output_csv = r"C:\Users\deverej\Presidio\pii_scan_test06.2025.csv"
log_file = r"C:\Users\deverej\Presidio\pii_scan_test06.2025.txt"

# Register HEIF/HEIC support with Pillow
register_heif_opener()

def extract_text_from_file(file_path):
    text_data = []
    try:
        if file_path.endswith(('.txt', '.csv', '.json')):
            with open(file_path, "r", encoding="utf-8") as file:
                text = file.read()
                pages = text.split('\f')  # Split text into pages using form feed
                for page_number, page_text in enumerate(pages, start=1):
                    text_data.append((page_text, page_number))
        elif file_path.endswith('.pdf'):
            reader = PdfReader(file_path)
            for page_number, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ''
                text_data.append((text, page_number))
        elif file_path.endswith('.docx'):
            doc = Document(file_path)
            for page_number, paragraph in enumerate(doc.paragraphs, start=1):
                text_data.append((paragraph.text, page_number))
        elif file_path.endswith('.doc'):
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(file_path)
            for page_number in range(1, doc.ComputeStatistics(2) + 1):  # 2 = wdStatisticPages
                text = doc.Range(doc.GoTo(1, 1, page_number).Start, doc.GoTo(1, 1, page_number + 1).Start).Text
                text_data.append((text, page_number))
            doc.Close()
            word.Quit()
        elif file_path.endswith('.pptx'):
            prs = Presentation(file_path)
            for slide_number, slide in enumerate(prs.slides, start=1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                text_data.append((slide_text, slide_number))
        elif file_path.endswith('.ppt'):
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt.Presentations.Open(file_path, WithWindow=False)
            for slide_number, slide in enumerate(presentation.Slides, start=1):
                slide_text = ""
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        slide_text += shape.TextFrame.TextRange.Text + "\n"
                text_data.append((slide_text, slide_number))
            presentation.Close()
            ppt.Quit()
        elif file_path.endswith(('.xlsx', '.xls')):
            try:
                df = pd.read_excel(file_path, sheet_name=None)
                for sheet_name, sheet in df.items():
                    text_data.append((sheet.to_string(index=False), sheet_name))
            except Exception as e:
                log_error(file_path, f"Error reading Excel: {e}")
        elif file_path.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.gif', '.heic')):
            try:
                with Image.open(file_path) as img:
                    text = pytesseract.image_to_string(img, lang='eng')  # Specify language for OCR
                    text_data.append((text, 'N/A'))  # HEIC files are single-page, so page number is 'N/A'
            except Exception as e:
                log_error(file_path, f"Error processing image: {e}")
        elif file_path.endswith('.msg'):
            msg = extract_msg.Message(file_path)
            text_data.append((msg.body, 'N/A'))
        elif file_path.endswith('.eml'):
            with open(file_path, 'rb') as f:
                msg = BytesParser(policy=policy.default).parse(f)
                text = msg.get_body(preferencelist=('plain')).get_content()
                text_data.append((text, 'N/A'))
        elif file_path.endswith('.zip'):
            with tempfile.TemporaryDirectory() as tmpdir:
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(tmpdir)
                    for root, dirs, files in os.walk(tmpdir):
                        for name in files:
                            extracted_text = extract_text_from_file(os.path.join(root, name))
                            text_data.extend(extracted_text)
        elif file_path.endswith(('.tiff', '.tif')):
            try:
                with Image.open(file_path) as img:
                    for page_number in range(getattr(img, "n_frames", 1)):
                        img.seek(page_number)
                        text = pytesseract.image_to_string(img)
                        text_data.append((text, page_number + 1))  # Store page number (1-based index)
            except Exception as e:
                log_error(file_path, f"Error processing TIFF image: {e}")
        else:
            log_error(file_path, "Unsupported file type")
    except Exception as e:
        log_error(file_path, f"Error extracting text: {e}")
    return text_data

def log_error(file_path, message):
    with open(log_file, "a", encoding="utf-8") as log:
        log.write(f"{file_path}: {message}\n")

def split_text(text, max_length=1000000):
    return [text[i:i + max_length] for i in range(0, len(text), max_length)]

# Scan a folder for PII and write results
def scan_folder_for_pii(folder_path):
    os.makedirs(os.path.dirname(output_csv), exist_ok=True)
    with open(output_csv, "w", newline='', encoding="utf-8") as csvfile:
        writer = csv.DictWriter(csvfile, fieldnames=['file_path', 'page_number', 'pii_entity', 'start', 'end', 'score', 'pii_value'])
        writer.writeheader()
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                print(f"Processing: {file_path}")
                log_error(file_path, "Processing")
                text_data = extract_text_from_file(file_path)
                for text, page_number in text_data:
                    if text:
                        # Split the text into smaller chunks
                        text_chunks = split_text(text, max_length=1000000)
                        for chunk in text_chunks:
                            results = analyzer.analyze(
                                text=chunk,
                                entities=["CREDIT_CARD", "PHONE_NUMBER", "IBAN_CODE", "US_BANK_NUMBER",
                                          "US_DRIVER_LICENSE", "US_ITIN", "US_PASSPORT", "US_SSN"],
                                language="en"
                            )
                            for result in results:
                                pii_value = text[result.start:result.end]
                                writer.writerow({
                                    'file_path': file_path,
                                    'page_number': page_number if page_number else 'N/A',
                                    'pii_entity': result.entity_type,
                                    'start': result.start,
                                    'end': result.end,
                                    'score': result.score,
                                    'pii_value': pii_value
                                })
                                print(f"PII found: {result.entity_type} ({pii_value}) in {file_path} on page {page_number if page_number else 'N/A'}")
                    else:
                        log_error(file_path, "No extractable text")

# Run the scan
scan_folder_for_pii(folder_path)